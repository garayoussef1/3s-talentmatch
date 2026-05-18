"""
Génère la présentation PowerPoint de soutenance — 3S TalentMatch
Run : python scripts/generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── Palette de couleurs ──────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # fond sombre
BLUE        = RGBColor(0x1E, 0x5C, 0xAA)   # bleu principal
ACCENT      = RGBColor(0xE8, 0x87, 0x10)   # orange 3S
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xFA)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY    = RGBColor(0x6B, 0x72, 0x80)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
RED_SOFT    = RGBColor(0xEF, 0x44, 0x44)

LOGO_3S_PATH    = r"c:\Users\youssef\Desktop\3s-talentmatch\frontend\src\assets\logo_3s.png"
LOGO_ESPRIT_PATH = r"C:\Users\youssef\Desktop\3s-talentmatch\esprit.png"
LOGO_PATH        = LOGO_3S_PATH  # alias pour compatibilité

W = Inches(13.33)   # largeur slide widescreen
H = Inches(7.5)     # hauteur slide widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # complètement vide


# ════════════════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════════════════

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_name="Calibri", font_size=14, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multi_para(slide, lines, x, y, w, h,
                   font_name="Calibri", font_size=13, color=DARK_TEXT,
                   line_spacing=None, bold_first=False):
    """Ajoute un bloc de texte multi-lignes."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txBox


def add_logo(slide, path, x, y, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, height=h)


def add_header_bar(slide, title, subtitle=None):
    """Barre de titre en haut (bleu foncé), avec trait orange + les deux logos."""
    add_rect(slide, 0, 0, W, Inches(1.45), NAVY)
    add_rect(slide, 0, Inches(1.45), W, Pt(5), ACCENT)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.22), Inches(9.5), Inches(0.75),
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.95), Inches(9.5), Inches(0.45),
                     font_size=14, color=RGBColor(0xBB, 0xCC, 0xEE))
    # Logo ESPRIT à droite
    add_logo(slide, LOGO_ESPRIT_PATH, Inches(11.55), Inches(0.08), Inches(1.35))
    # Logo 3S juste à sa gauche
    add_logo(slide, LOGO_3S_PATH, Inches(10.2), Inches(0.15), Inches(1.1))


TOTAL_SLIDES = 15  # mis à jour après génération

def add_footer(slide, slide_num=None,
               text="Youssef GARA  |  Stage 3S  |  Enc. academique : Mme Loujayne Bouzrati  |  Enc. societe : Mme Sirine Nasri  |  2026"):
    add_rect(slide, 0, Inches(7.15), W, Inches(0.35), NAVY)
    # Texte à gauche
    add_text_box(slide, text, Inches(0.4), Inches(7.17), Inches(11.0), Inches(0.3),
                 font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC))
    # Numéro à droite
    if slide_num is not None:
        num_txt = f"{slide_num} / {TOTAL_SLIDES}"
        add_text_box(slide, num_txt, Inches(11.6), Inches(7.17), Inches(1.5), Inches(0.3),
                     font_size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def bullet_box(slide, items, x, y, w, h,
               font_size=14, color=DARK_TEXT, bullet="  ›  "):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet + item
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def section_title(slide, text, x, y, w=Inches(5)):
    """Titre de section avec fond bleu."""
    box = add_rect(slide, x, y, w, Pt(32), BLUE)
    add_text_box(slide, text, x + Inches(0.15), y + Pt(4), w - Inches(0.3), Pt(28),
                 font_size=14, bold=True, color=WHITE)


def card(slide, title, lines, x, y, w=Inches(5.8), h=Inches(1.9),
         bg=LIGHT_GRAY, title_color=BLUE, text_color=DARK_TEXT, font_size=12):
    add_rect(slide, x, y, w, h, bg)
    add_rect(slide, x, y, w, Pt(6), title_color)
    add_text_box(slide, title, x + Inches(0.15), y + Pt(12), w, Pt(24),
                 font_size=13, bold=True, color=title_color)
    bullet_box(slide, lines, x + Inches(0.1), y + Pt(36), w - Inches(0.2),
               h - Pt(42), font_size=font_size, color=text_color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Couverture
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Fond dégradé simulé : deux rectangles
add_rect(slide, 0, 0, W, H, NAVY)
add_rect(slide, 0, Inches(5.2), W, Inches(2.3), BLUE)
add_rect(slide, 0, Inches(5.1), W, Pt(6), ACCENT)

# Logo 3S en haut à gauche
add_logo(slide, LOGO_3S_PATH, Inches(0.5), Inches(0.3), Inches(1.5))
# Logo ESPRIT en haut à droite
add_logo(slide, LOGO_ESPRIT_PATH, Inches(11.3), Inches(0.15), Inches(1.7))

# Titre principal
add_text_box(slide, "TalentMatch",
             Inches(2.5), Inches(1.0), Inches(9), Inches(1.4),
             font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Sous-titre
add_text_box(slide,
             "Plateforme intelligente de recrutement par IA sémantique",
             Inches(1.5), Inches(2.4), Inches(10.3), Inches(0.7),
             font_size=20, italic=True,
             color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

# Trait orange
add_rect(slide, Inches(4.5), Inches(3.15), Inches(4.3), Pt(4), ACCENT)

# Infos
add_text_box(slide, "Youssef GARA",
             Inches(0), Inches(3.5), W, Inches(0.5),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, "Licence en Genie Logiciel  —  ESPRIT  |  Stage PFE 6 mois",
             Inches(0), Inches(4.0), W, Inches(0.45),
             font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_text_box(slide, "Enc. académique : Mme Loujayne Bouzrati   •   Enc. société : Mme Sirine Nasri",
             Inches(0), Inches(4.45), W, Inches(0.45),
             font_size=13, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_text_box(slide, "Fevrier 2026 — Aout 2026",
             Inches(0), Inches(5.8), W, Inches(0.45),
             font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

# Badge "IA"
badge = add_rect(slide, Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.7), ACCENT)
add_text_box(slide, "🤖 IA / NLP",
             Inches(11.5), Inches(0.43), Inches(1.5), Inches(0.5),
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Plan de la présentation
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Plan de la présentation")
add_footer(slide, 2)

# (label, sous-titre, slide_num, couleur_badge)
sommaire_items = [
    ("Introduction",                  "Presentation de 3S, contexte, chiffres cles",  3,  BLUE),
    ("Cadre du projet",               "Problematique, etude existant, solution",       4,  BLUE),
    ("Analyse des besoins",           "Besoins fonctionnels & non fonctionnels",       6,  BLUE),
    ("Architectures",                 "Stack 3 tiers : React / FastAPI / PostgreSQL",  8,  BLUE),
    ("Methodologie",                  "Agile Scrum — 6 Sprints d'1 mois",              10, ACCENT),
    ("Avancement actuel",             "Fonctionnalites livrees + demonstration",       11, ACCENT),
    ("Conclusion & Perspectives",     "Bilan technique, acquis, evolutions futures",   13, ACCENT),
]

cols_items = [sommaire_items[:4], sommaire_items[4:]]
# Note: 7 items -> col gauche 4 items, col droite 3 items
col_x = [Inches(0.4), Inches(6.85)]

for col_i, col_list in enumerate(cols_items):
    x = col_x[col_i]
    color = BLUE if col_i == 0 else ACCENT
    for row_i, (label, sub, pg, _) in enumerate(col_list):
        y_pos = Inches(1.68) + row_i * Inches(1.3)
        # Carte blanche
        add_rect(slide, x, y_pos, Inches(6.2), Inches(1.15), WHITE)
        # Barre couleur à gauche
        add_rect(slide, x, y_pos, Pt(6), Inches(1.15), color)
        # Badge numéro en haut à droite de la carte
        add_rect(slide, x + Inches(5.55), y_pos + Pt(6), Inches(0.55), Inches(0.45), color)
        add_text_box(slide, f"{pg:02d}", x + Inches(5.55), y_pos + Pt(8),
                     Inches(0.55), Pt(26), font_size=13, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
        # Titre
        add_text_box(slide, label, x + Inches(0.22), y_pos + Pt(12),
                     Inches(5.2), Pt(26), font_size=15, bold=True, color=DARK_TEXT)
        # Sous-titre
        add_text_box(slide, sub, x + Inches(0.22), y_pos + Pt(42),
                     Inches(5.2), Pt(26), font_size=11, color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction & Contexte
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Introduction",
               "Presentation de 3S Tunisie — Contexte & enjeux du projet")
add_footer(slide, 3)

# ── Chiffres clés 3S (bandeau haut) ────────────────────────────────────────
kpi_items = [
    ("1988", "Annee de fondation"),
    ("35+", "Ans d'expertise IT"),
    ("300+", "Employes"),
    ("450+", "Certifications"),
    ("2 000+", "Clients satisfaits"),
]
for i, (val, lbl) in enumerate(kpi_items):
    x = Inches(0.3) + i * Inches(2.6)
    add_rect(slide, x, Inches(1.62), Inches(2.45), Inches(0.95), NAVY)
    add_text_box(slide, val, x, Inches(1.65), Inches(2.45), Pt(28),
                 font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, lbl, x, Inches(1.98), Inches(2.45), Pt(20),
                 font_size=10, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

# ── Carte gauche — Présentation 3S ──────────────────────────────────────────
add_rect(slide, Inches(0.4), Inches(2.75), Inches(5.9), Inches(2.6), WHITE)
add_rect(slide, Inches(0.4), Inches(2.75), Pt(6), Inches(2.6), BLUE)
add_text_box(slide, "L'entreprise — 3S Tunisie",
             Inches(0.6), Inches(2.82), Inches(5.5), Pt(26),
             font_size=14, bold=True, color=BLUE)
bullet_box(slide, [
    "Fondee en 1988, leader tunisien de l'integration IT",
    "Specialiste : infrastructures reseaux, cybersecurite,",
    "  data centers, cloud, espaces de travail numeriques",
    "ISO 9001:2015 (Qualite) + ISO 27001:2013 (Securite)",
    "1er integrateur IT tunisien certifie ISO 27001",
    "Expansion : filiale a Abidjan (Cote d'Ivoire) en 2019",
    "Partenaire strategique de la transformation numerique",
], Inches(0.6), Inches(3.18), Inches(5.5), Inches(2.0),
   font_size=12, color=DARK_TEXT)

# ── Carte droite — Contexte du stage ────────────────────────────────────────
add_rect(slide, Inches(7.0), Inches(2.75), Inches(5.9), Inches(2.6), WHITE)
add_rect(slide, Inches(7.0), Inches(2.75), Pt(6), Inches(2.6), ACCENT)
add_text_box(slide, "Contexte du stage",
             Inches(7.2), Inches(2.82), Inches(5.5), Pt(26),
             font_size=14, bold=True, color=ACCENT)
bullet_box(slide, [
    "Stage de fin d'etudes — Licence Genie Logiciel, ESPRIT",
    "Duree : 6 mois (Fevrier 2026 -> Aout 2026)",
    "Enc. academique : Mme Loujayne Bouzrati (ESPRIT)",
    "Enc. societe : Mme Sirine Nasri (3S)",
    "Mission : concevoir une plateforme IA de recrutement",
    "Objectif : automatiser le matching CV <-> offres d'emploi",
], Inches(7.2), Inches(3.18), Inches(5.5), Inches(2.0),
   font_size=12, color=DARK_TEXT)

# ── Bandeau bas : besoin identifie ──────────────────────────────────────────
add_rect(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.2), NAVY)
add_text_box(slide, "Besoin identifie chez 3S",
             Inches(0.6), Inches(5.58), Inches(5), Pt(24),
             font_size=13, bold=True, color=ACCENT)
bullet_box(slide, [
    "Le tri manuel des CVs est chronophage, subjectif et non scalable",
    "Besoin d'un outil interne intelligent : score automatique, multi-domaines, expliquable",
], Inches(0.6), Inches(5.88), Inches(12.0), Inches(0.75),
   font_size=12, color=WHITE, bullet="  ->  ")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Problématique & Étude de l'existant
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Cadre du projet — Problematique & Etude de l'existant",
               "Pourquoi ce projet ? Que font les autres ? Quelle est notre reponse ?")
add_footer(slide, 4)

# Problématique — colonne gauche
add_rect(slide, Inches(0.4), Inches(1.65), Inches(4.5), Inches(5.1), WHITE)
add_rect(slide, Inches(0.4), Inches(1.65), Pt(5), Inches(5.1), RED_SOFT)
add_text_box(slide, "❌  Problème actuel",
             Inches(0.6), Inches(1.75), Inches(4.2), Inches(0.4),
             font_size=14, bold=True, color=RED_SOFT)
bullet_box(slide, [
    "Tri manuel des CVs : lent et subjectif",
    "Mots-clés insuffisants : manque la sémantique",
    "Aucun retour au candidat sur son profil",
    "Pas d'apprentissage des décisions passées",
    "Outils existants trop génériques ou trop coûteux",
], Inches(0.6), Inches(2.2), Inches(4.2), Inches(4.2),
   font_size=13, color=DARK_TEXT)

# Tableau comparatif — colonne droite
add_rect(slide, Inches(5.2), Inches(1.65), Inches(7.7), Inches(5.1), WHITE)
add_rect(slide, Inches(5.2), Inches(1.65), Inches(7.7), Pt(5), BLUE)
add_text_box(slide, "🔍  Comparaison des solutions existantes",
             Inches(5.4), Inches(1.75), Inches(7.3), Inches(0.4),
             font_size=14, bold=True, color=BLUE)

headers = ["Solution", "Matching sémantique", "NLP Français", "Coût", "Personnalisable"]
rows = [
    ["LinkedIn Recruiter",  "✅", "Partiel", "💰💰💰", "❌"],
    ["Workday / Taleo",     "Partiel", "❌", "💰💰💰", "Partiel"],
    ["Jobology / Flatchr",  "Partiel", "✅", "💰💰", "Partiel"],
    ["CVthèque Excel",      "❌", "❌", "Gratuit", "❌"],
    ["3S TalentMatch ✨",   "✅ BERT", "✅ spaCy FR", "Gratuit", "✅ 100%"],
]

col_w = [Inches(2.2), Inches(1.5), Inches(1.3), Inches(1.1), Inches(1.5)]
col_x = [Inches(5.25), Inches(7.45), Inches(8.95), Inches(10.25), Inches(11.35)]
y_start = Inches(2.2)
row_h = Inches(0.58)

for c_i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, cx, y_start, cw, row_h, BLUE)
    add_text_box(slide, hdr, cx + Pt(4), y_start + Pt(8), cw, row_h,
                 font_size=11, bold=True, color=WHITE)

for r_i, row in enumerate(rows):
    y = y_start + (r_i + 1) * row_h
    bg = RGBColor(0xE8, 0xF0, 0xFE) if row[0].endswith("✨") else (
         LIGHT_GRAY if r_i % 2 == 0 else WHITE)
    for c_i, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, cx, y, cw, row_h, bg)
        clr = GREEN if row[0].endswith("✨") else DARK_TEXT
        add_text_box(slide, cell, cx + Pt(4), y + Pt(6), cw, row_h,
                     font_size=11, color=clr)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Solution proposée
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Solution proposée — 3S TalentMatch",
               "Une plateforme end-to-end : upload CV → score IA → décision recruteur")
add_footer(slide, 5)

# 3 rôles
roles = [
    ("👤 Candidat", BLUE, [
        "S'inscrit & dépose son CV (PDF/DOCX)",
        "Postule aux offres disponibles",
        "Reçoit feedback : accepté / refusé",
        "Consulte son score de matching",
    ]),
    ("🧑‍💼 Recruteur", ACCENT, [
        "Crée et publie des offres d'emploi",
        "Reçoit les candidats classés par l'IA",
        "Accepte / refuse en un clic",
        "Suit les candidatures en temps réel",
    ]),
    ("⚙️  Administrateur", NAVY, [
        "Gère les utilisateurs & recruteurs",
        "Consulte les logs d'accès (RGPD)",
        "Surveille les métriques IA",
        "Vue globale du système",
    ]),
]

for i, (title, color, items) in enumerate(roles):
    x = Inches(0.4) + i * Inches(4.3)
    add_rect(slide, x, Inches(1.65), Inches(4.1), Inches(2.6), WHITE)
    add_rect(slide, x, Inches(1.65), Inches(4.1), Pt(6), color)
    add_text_box(slide, title, x + Inches(0.15), Inches(1.75),
                 Inches(3.8), Inches(0.45), font_size=14, bold=True, color=color)
    bullet_box(slide, items, x + Inches(0.1), Inches(2.25),
               Inches(3.9), Inches(1.9), font_size=12, color=DARK_TEXT)

# Flux principal
add_rect(slide, Inches(0.4), Inches(4.45), Inches(12.5), Inches(2.65), WHITE)
add_rect(slide, Inches(0.4), Inches(4.45), Inches(12.5), Pt(5), GREEN)
add_text_box(slide, "🔄  Flux automatisé",
             Inches(0.6), Inches(4.55), Inches(5), Inches(0.4),
             font_size=14, bold=True, color=GREEN)

steps = [
    ("1. Upload CV", "PDF/DOCX/Image"),
    ("2. Extraction", "Texte brut"),
    ("3. NLP Parser", "Skills, Exp, Formation"),
    ("4. BERT Matching", "Similarité sémantique"),
    ("5. Score + Décision", "0–100% + label"),
    ("6. Notification", "Email + in-app"),
]
for i, (step, sub) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.15)
    add_rect(slide, x, Inches(5.0), Inches(2.0), Inches(0.95), LIGHT_GRAY)
    add_text_box(slide, step, x + Pt(6), Inches(5.05), Inches(1.9), Pt(22),
                 font_size=11, bold=True, color=BLUE)
    add_text_box(slide, sub, x + Pt(6), Inches(5.3), Inches(1.9), Pt(22),
                 font_size=10, color=MID_GRAY)
    if i < 5:
        add_text_box(slide, "→", x + Inches(2.0), Inches(5.1), Pt(30), Pt(30),
                     font_size=18, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Analyse des besoins fonctionnels
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Analyse des besoins fonctionnels")
add_footer(slide, 6)

func_groups = [
    ("👤 Candidat", BLUE, Inches(0.4), [
        "S'inscrire par email ou OAuth (Google / LinkedIn)",
        "Uploader son CV (PDF, DOCX, image)",
        "Consulter les offres d'emploi disponibles",
        "Postuler à une offre et suivre son statut",
        "Recevoir une notification (email + in-app)",
    ]),
    ("🧑‍💼 Recruteur", ACCENT, Inches(5.1), [
        "Créer, modifier et supprimer des offres",
        "Lancer le matching IA pour une offre",
        "Voir les candidats classés par score",
        "Changer le statut d'une candidature",
        "Recevoir une alerte à chaque nouveau CV",
    ]),
    ("⚙️  Admin", NAVY, Inches(9.8), [
        "Gérer tous les utilisateurs",
        "Créer des comptes recruteurs",
        "Consulter les logs RGPD",
        "Accéder au dashboard global",
    ]),
]

for title, color, x, items in func_groups:
    w = Inches(4.55) if x < Inches(9) else Inches(3.2)
    add_rect(slide, x, Inches(1.65), w, Inches(5.1), WHITE)
    add_rect(slide, x, Inches(1.65), w, Pt(5), color)
    add_text_box(slide, title, x + Inches(0.15), Inches(1.75), w, Pt(28),
                 font_size=14, bold=True, color=color)
    bullet_box(slide, items, x + Inches(0.1), Inches(2.15), w - Inches(0.2),
               Inches(4.4), font_size=12.5, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Besoins non fonctionnels
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Analyse des besoins non fonctionnels")
add_footer(slide, 7)

nf_items = [
    ("🔐 Sécurité", BLUE, [
        "JWT + refresh token",
        "OAuth Google & LinkedIn",
        "Hashage bcrypt des mots de passe",
        "HTTPS en production",
    ]),
    ("⚡ Performance", GREEN, [
        "BERT chargé au démarrage (cold start évité)",
        "Réponse matching < 2 sec par candidat",
        "PostgreSQL + index sur clés étrangères",
    ]),
    ("📋 RGPD", ACCENT, [
        "Logs d'accès traçables (user, action, date)",
        "Suppression de CV sur demande",
        "Données localisées (Tunisie)",
    ]),
    ("🔧 Maintenabilité", NAVY, [
        "Architecture modulaire (services séparés)",
        "Moteurs IA interchangeables (heuristique / BERT)",
        "API REST documentée (Swagger /docs)",
    ]),
    ("📱 Utilisabilité", RGBColor(0x7C, 0x3A, 0xED), [
        "Interface responsive (mobile + desktop)",
        "Feedback temps réel (notifications cloche)",
        "Messages d'erreur explicites",
    ]),
    ("🚀 Extensibilité", RGBColor(0x06, 0x95, 0x6E), [
        "Nouveaux moteurs IA en 1 fichier",
        "Nouveaux rôles sans refonte",
        "API versionnée (/api/v1/...)",
    ]),
]

cols = 3
for i, (title, color, items) in enumerate(nf_items):
    row, col = divmod(i, cols)
    x = Inches(0.4) + col * Inches(4.3)
    y = Inches(1.7) + row * Inches(2.55)
    add_rect(slide, x, y, Inches(4.1), Inches(2.35), WHITE)
    add_rect(slide, x, y, Pt(6), Inches(2.35), color)
    add_text_box(slide, title, x + Inches(0.2), y + Pt(10), Inches(3.8), Pt(28),
                 font_size=13, bold=True, color=color)
    bullet_box(slide, items, x + Inches(0.2), y + Pt(40),
               Inches(3.8), Inches(2.35) - Pt(46), font_size=12, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Architecture technique
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Architecture technique",
               "Stack 3 tiers : React → FastAPI → PostgreSQL")
add_footer(slide, 8)

layers = [
    ("🖥️  Frontend", BLUE, "React + Vite", [
        "Composants JSX",
        "React Router v6",
        "AuthContext (JWT)",
        "Axios (REST calls)",
    ]),
    ("⚙️  Backend", ACCENT, "FastAPI (Python)", [
        "API REST (8 modules)",
        "SQLAlchemy ORM",
        "JWT / OAuth",
        "Pipeline NLP + BERT",
    ]),
    ("🗄️  Base de données", GREEN, "PostgreSQL + Alembic", [
        "Users, CVs, JobOffers",
        "Candidates, Matches",
        "Notifications, Logs",
        "Migrations versionnées",
    ]),
    ("🤖  IA / NLP", NAVY, "BERT + spaCy", [
        "TalentMatch-BERT v2.0",
        "MLP 5 dimensions",
        "spaCy fr_core_news_md",
        "14 domaines métier",
    ]),
]

for i, (title, color, subtitle, items) in enumerate(layers):
    x = Inches(0.4) + i * Inches(3.2)
    add_rect(slide, x, Inches(1.7), Inches(3.0), Inches(4.9), WHITE)
    add_rect(slide, x, Inches(1.7), Inches(3.0), Pt(5), color)
    # Header
    add_rect(slide, x, Inches(1.7), Inches(3.0), Inches(0.9), color)
    add_text_box(slide, title, x + Pt(8), Inches(1.73), Inches(2.8), Pt(26),
                 font_size=13, bold=True, color=WHITE)
    add_text_box(slide, subtitle, x + Pt(8), Inches(2.02), Inches(2.8), Pt(22),
                 font_size=11, color=RGBColor(0xDD, 0xEE, 0xFF))
    # Contenu
    bullet_box(slide, items, x + Pt(8), Inches(2.65), Inches(2.8), Inches(3.7),
               font_size=12.5, color=DARK_TEXT, bullet="• ")

# Flèches entre couches
for i in range(3):
    x = Inches(3.42) + i * Inches(3.2)
    add_text_box(slide, "⟷", x, Inches(3.85), Inches(0.18), Inches(0.5),
                 font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Note bas
add_rect(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.35), NAVY)
add_text_box(slide,
             "Communication : CORS sécurisé  •  JWT Bearer token  •  SMTP (email)  •  Alembic (migrations)",
             Inches(0.6), Inches(6.77), Inches(12), Inches(0.3),
             font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Pipeline IA & NLP
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Moteur de Matching IA — Pipeline complet",
               "De l'octet brut au score de pertinence")
add_footer(slide, 9)

pipeline_steps = [
    ("📄 CV brut", "PDF / DOCX / Image",    BLUE),
    ("📝 Extraction", "PyMuPDF · python-docx · OCR Tesseract", BLUE),
    ("🧠 NLP Parser", "spaCy FR · extracteurs custom\nSkills · Expériences · Formation · Contact", ACCENT),
    ("🤖 BERT Scorer", "TalentMatch-BERT v2.0\n(fine-tuné sur 4 966 CVs réels)", NAVY),
    ("📊 MLP 5D", "[sem_sim · skill_rate · exp_score\nformation · appreciated_rate]", NAVY),
    ("✅ Score", "0–100%\nExcellent / Bon / À évaluer / Non adapté", GREEN),
]

box_w = Inches(1.95)
box_h = Inches(1.35)
for i, (title, sub, color) in enumerate(pipeline_steps):
    x = Inches(0.35) + i * Inches(2.17)
    y = Inches(1.75)
    add_rect(slide, x, y, box_w, box_h, color)
    add_text_box(slide, title, x + Pt(6), y + Pt(8), box_w - Pt(8), Pt(24),
                 font_size=12, bold=True, color=WHITE)
    add_text_box(slide, sub, x + Pt(6), y + Pt(32), box_w - Pt(8), box_h - Pt(38),
                 font_size=9, color=RGBColor(0xCC, 0xDD, 0xEE))
    if i < 5:
        add_text_box(slide, "▶", x + box_w + Pt(4), y + Pt(35), Pt(28), Pt(28),
                     font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Détail BERT
add_rect(slide, Inches(0.4), Inches(3.35), Inches(12.5), Inches(3.4), WHITE)
add_rect(slide, Inches(0.4), Inches(3.35), Inches(12.5), Pt(4), NAVY)
add_text_box(slide, "🔬  Détail du scoring BERT v2.0",
             Inches(0.6), Inches(3.42), Inches(8), Inches(0.4),
             font_size=14, bold=True, color=NAVY)

bert_details = [
    ("📐 Similarité sémantique\n(sem_sim)", "Embeddings 384D\nMiniLM multilingual\nCV ↔ Offre cosine"),
    ("🎯 Taux skills requis\n(skill_rate)", "RapidFuzz ≥ 80%\nFuzzy matching FR+EN\n4 niveaux de vérif."),
    ("📅 Score expérience\n(exp_score)", "Années extraites NLP\nNormalisé vs offre\nPénalité si <50%"),
    ("🎓 Niveau formation\n(form_score)", "Bac+2 → Bac+5+\nMatching vs exigences\nExtracteur NLP dédié"),
    ("⭐ Skills appréciées\n(appreciated)", "Bonus soft skills\nSkills complémentaires\n0–30% bonus"),
]
for i, (title, body) in enumerate(bert_details):
    x = Inches(0.5) + i * Inches(2.5)
    add_rect(slide, x, Inches(3.85), Inches(2.35), Inches(2.65), LIGHT_GRAY)
    add_text_box(slide, title, x + Pt(6), Inches(3.92), Inches(2.2), Pt(36),
                 font_size=11, bold=True, color=NAVY)
    add_text_box(slide, body, x + Pt(6), Inches(4.45), Inches(2.2), Inches(1.9),
                 font_size=10.5, color=MID_GRAY)

# Domain caps
add_rect(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.55), NAVY)
add_text_box(slide,
             "🔒  Caps par domaine :  même domaine → pas de cap  |  domaine adjacent (IT↔Ingénierie) → 62%  |  domaine différent → 44%  |  Floor minimum → 20%",
             Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             font_size=11, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Méthodologie Agile / Sprints
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Methodologie — Approche Agile Scrum",
               "Stage PFE 6 mois  |  6 Sprints d'1 mois  |  Fevrier 2026 - Aout 2026")
add_footer(slide, 10)

# ── Barre de progression ─────────────────────────────────────────────────────
TL_X = Inches(0.38)
TL_W = Inches(12.55)
TL_Y = Inches(1.72)
# Fond gris (total)
add_rect(slide, TL_X, TL_Y, TL_W, Pt(6), RGBColor(0xCC, 0xD5, 0xE0))
# Progression : Sprints 1-3 termines = 3/6 = 50%
add_rect(slide, TL_X, TL_Y, Inches(6.27), Pt(6), BLUE)
# Sprint 4 en cours = +16% = jusqu'a ~60%
add_rect(slide, TL_X + Inches(6.27), TL_Y, Inches(2.09), Pt(6), ACCENT)

# Mois sur la timeline
months = [
    ("Fev",  Inches(0.38)),
    ("Mar",  Inches(2.47)),
    ("Avr",  Inches(4.56)),
    ("Mai",  Inches(6.64)),
    ("Juin", Inches(8.72)),
    ("Juil", Inches(10.8)),
    ("Aout", Inches(12.55)),
]
for lbl, x in months:
    add_rect(slide, x, Inches(1.64), Pt(2), Pt(20), RGBColor(0x88, 0x99, 0xAA))
    add_text_box(slide, lbl, x - Inches(0.28), Inches(1.86),
                 Inches(0.9), Pt(16), font_size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Marqueur "1ere restitution" sur Mai 2026 (entre Sprint 3 et 4)
mk_x = Inches(6.6)
add_rect(slide, mk_x, Inches(1.57), Pt(3), Inches(0.5), ACCENT)
add_rect(slide, mk_x - Inches(0.72), Inches(1.3), Inches(1.55), Pt(22), ACCENT)
add_text_box(slide, "1ere restitution",
             mk_x - Inches(0.7), Inches(1.31),
             Inches(1.52), Pt(18),
             font_size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── 6 Sprints ────────────────────────────────────────────────────────────────
sprints = [
    {
        "title":    "Sprint 1",
        "mois":     "Fevrier 2026",
        "color":    BLUE,
        "status":   "TERMINE",
        "status_c": GREEN,
        "tasks": [
            "Cahier des charges",
            "Conception UML",
            "Architecture technique",
            "Setup environnement",
            "(backend, BDD, frontend)",
        ],
    },
    {
        "title":    "Sprint 2",
        "mois":     "Mars 2026",
        "color":    RGBColor(0x1A, 0x6F, 0xA0),
        "status":   "TERMINE",
        "status_c": GREEN,
        "tasks": [
            "Upload CV (PDF/DOCX)",
            "Extraction texte",
            "Integration OCR",
            "(EasyOCR)",
            "Stockage BDD",
        ],
    },
    {
        "title":    "Sprint 3",
        "mois":     "Avril 2026",
        "color":    RGBColor(0x0D, 0x4A, 0x70),
        "status":   "TERMINE",
        "status_c": GREEN,
        "tasks": [
            "Pipeline NLP spaCy",
            "Extraction skills",
            "Extraction experiences",
            "Extraction formations",
            "Tests unitaires",
        ],
    },
    {
        "title":    "Sprint 4",
        "mois":     "Mai 2026",
        "color":    ACCENT,
        "status":   "EN COURS",
        "status_c": ACCENT,
        "tasks": [
            "BERT / CamemBERT",
            "Matching candidat-offre",
            "Score IA (0-100%)",
            "Dashboard RH",
            "Optimisation perfs",
        ],
    },
    {
        "title":    "Sprint 5",
        "mois":     "Juin 2026",
        "color":    MID_GRAY,
        "status":   "A VENIR",
        "status_c": MID_GRAY,
        "tasks": [
            "Securite & Auth",
            "Roles & RGPD",
            "Tests integration",
            "Ameliorations UX",
            "Notifications",
        ],
    },
    {
        "title":    "Sprint 6",
        "mois":     "Juil - Aout 2026",
        "color":    RGBColor(0x55, 0x5F, 0x6A),
        "status":   "A VENIR",
        "status_c": MID_GRAY,
        "tasks": [
            "Validation finale",
            "Documentation",
            "Deploiement",
            "Preparation soutenance",
        ],
    },
]

# Dimensions pour 6 cartes
CARD_W = Inches(2.02)
CARD_H = Inches(4.62)
GAP    = Inches(0.085)

for i, sprint in enumerate(sprints):
    x = Inches(0.38) + i * (CARD_W + GAP)
    color     = sprint["color"]
    is_future = (sprint["status"] == "A VENIR")
    is_active = (sprint["status"] == "EN COURS")

    # Ombre
    add_rect(slide, x + Pt(3), Inches(2.12) + Pt(3), CARD_W, CARD_H,
             RGBColor(0xBB, 0xBB, 0xBB))

    # Fond carte
    card_bg = RGBColor(0xF2, 0xF2, 0xF2) if is_future else WHITE
    add_rect(slide, x, Inches(2.12), CARD_W, CARD_H, card_bg)

    # Bordure gauche colorée
    add_rect(slide, x, Inches(2.12), Pt(5), CARD_H, color)

    # Header coloré
    header_h = Inches(0.65)
    add_rect(slide, x, Inches(2.12), CARD_W, header_h, color)

    # Titre
    add_text_box(slide, sprint["title"],
                 x + Pt(7), Inches(2.14),
                 Inches(1.3), Pt(24),
                 font_size=14, bold=True, color=WHITE)

    # Badge statut (petit, en haut à droite du header)
    sc = sprint["status_c"]
    badge_w = Inches(0.65)
    add_rect(slide, x + CARD_W - badge_w - Pt(4), Inches(2.22),
             badge_w, Pt(16), WHITE)
    lbl_color = sc if not is_active else ACCENT
    short_lbl = "OK" if sprint["status"] == "TERMINE" else ("..." if is_active else "TODO")
    add_text_box(slide, short_lbl,
                 x + CARD_W - badge_w - Pt(4), Inches(2.22),
                 badge_w, Pt(16),
                 font_size=8, bold=True, color=lbl_color, align=PP_ALIGN.CENTER)

    # Mois
    add_text_box(slide, sprint["mois"],
                 x + Pt(7), Inches(2.8),
                 CARD_W - Pt(10), Pt(18),
                 font_size=8, italic=True, color=MID_GRAY)

    # Séparateur
    add_rect(slide, x + Pt(6), Inches(3.02),
             CARD_W - Pt(10), Pt(2), LIGHT_GRAY)

    # Tâches
    txt_color = RGBColor(0xAA, 0xAA, 0xAA) if is_future else DARK_TEXT
    bullet_box(slide, sprint["tasks"],
               x + Pt(7), Inches(3.08),
               CARD_W - Pt(12), Inches(3.5),
               font_size=10.5, color=txt_color, bullet=" > ")

    # Flèche entre sprints
    if i < 5:
        arrow_x = x + CARD_W + Pt(1)
        add_text_box(slide, ">",
                     arrow_x, Inches(4.3),
                     GAP + Pt(4), Pt(24),
                     font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ── Note bas ─────────────────────────────────────────────────────────────────
add_rect(slide, Inches(0.38), Inches(6.97), Inches(12.55), Pt(16), NAVY)
add_text_box(slide,
             "Agile Scrum  |  Reviews regulieres avec encadrantes  |  Git  |  VS Code  |  PyCharm  |  Postman  |  Trello",
             Inches(0.55), Inches(6.985), Inches(12.2), Pt(14),
             font_size=9, color=RGBColor(0xAA, 0xBB, 0xCC))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Avancement actuel
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Avancement actuel — Fonctionnalites livrees")
add_footer(slide, 11)

features = [
    ("🔐 Authentification", GREEN, [
        "JWT Bearer + Refresh token",
        "OAuth Google & LinkedIn",
        "Vérification email (SMTP)",
        "3 rôles : candidat / recruteur / admin",
    ]),
    ("📄 Extraction CV", GREEN, [
        "PDF textuel (PyMuPDF)",
        "DOCX (python-docx)",
        "PDF scanné (OCR Tesseract)",
        "Formats : PDF, DOCX, PNG, JPG",
    ]),
    ("🧠 NLP Parser", GREEN, [
        "spaCy fr_core_news_md",
        "Extraction : compétences, expériences,\nformations, contact",
        "Support FR + EN",
        "Détection domaine métier (14 domaines)",
    ]),
    ("🤖 Matching IA", GREEN, [
        "BERT v2.0 fine-tuné (4 966 CVs)",
        "MLP 5D calibré (val_loss 0.00070)",
        "Score 0–100% + décision",
        "Caps par domaine métier",
    ]),
    ("🔔 Notifications", GREEN, [
        "Email candidat (accepté / refusé)",
        "Cloche in-app (temps réel)",
        "Dashboard statistiques recruteur",
        "Logs d'accès RGPD",
    ]),
    ("🖥️  Interface React", GREEN, [
        "Pages candidat, recruteur, admin",
        "Upload CV drag & drop",
        "Tableau de bord matching",
        "Gestion des offres",
    ]),
]

cols = 3
for i, (title, color, items) in enumerate(features):
    row, col = divmod(i, cols)
    x = Inches(0.4) + col * Inches(4.3)
    y = Inches(1.7) + row * Inches(2.55)
    add_rect(slide, x, y, Inches(4.1), Inches(2.35), WHITE)
    add_rect(slide, x, y, Pt(6), Inches(2.35), color)
    add_text_box(slide, title, x + Inches(0.2), y + Pt(10), Inches(3.8), Pt(26),
                 font_size=13, bold=True, color=color)
    bullet_box(slide, items, x + Inches(0.15), y + Pt(38),
               Inches(3.8), Inches(2.35) - Pt(44), font_size=11.5, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Démonstration (captures mockup texte)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Demonstration — Parcours type", "Upload CV -> Matching IA -> Decision recruteur")
add_footer(slide, 12)

steps_demo = [
    ("Étape 1", "📄 Candidat uploade son CV",
     "Format PDF/DOCX → extraction automatique du texte\n→ Parsing NLP : skills, expériences, formations extraits"),
    ("Étape 2", "🔍 Recruteur lance le matching",
     "POST /api/match/{offer_id}\n→ BERT calcule le score pour chaque candidat\n→ Classement automatique"),
    ("Étape 3", "📊 Résultats affichés",
     "Score 0–100% + décision\n✅ Excellent  🟢 Bon profil  🟡 À évaluer  🔴 Non adapté\n+ Skills matchées détaillées"),
    ("Étape 4", "✅ Recruteur décide",
     "Accepte ou refuse en un clic\n→ Email automatique au candidat\n→ Notification in-app mise à jour"),
]

for i, (step_num, step_title, step_body) in enumerate(steps_demo):
    row, col = divmod(i, 2)
    x = Inches(0.4) + col * Inches(6.5)
    y = Inches(1.7) + row * Inches(2.55)
    w, h = Inches(6.2), Inches(2.35)
    add_rect(slide, x, y, w, h, WHITE)
    # Badge étape
    add_rect(slide, x, y, Inches(1.1), h, NAVY)
    add_text_box(slide, step_num, x + Pt(4), y + Pt(30), Inches(1.0), Pt(28),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Contenu
    add_text_box(slide, step_title, x + Inches(1.15), y + Pt(14),
                 Inches(4.9), Pt(26), font_size=13, bold=True, color=NAVY)
    add_text_box(slide, step_body, x + Inches(1.15), y + Pt(42),
                 Inches(4.9), h - Pt(48), font_size=11.5, color=MID_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Conclusion")
add_footer(slide, 13)

# Ce qu'on a construit
add_rect(slide, Inches(0.4), Inches(1.7), Inches(8.5), Inches(5.1), WHITE)
add_rect(slide, Inches(0.4), Inches(1.7), Pt(5), Inches(5.1), GREEN)
add_text_box(slide, "✅  Ce que TalentMatch apporte",
             Inches(0.6), Inches(1.8), Inches(8), Pt(28), font_size=14, bold=True, color=GREEN)
bullet_box(slide, [
    "Un pipeline IA end-to-end : de l'upload CV au score automatique",
    "Un modèle BERT fine-tuné sur données réelles FR+EN (4 966 CVs, 24 domaines)",
    "Un matching sémantique qui va au-delà des mots-clés",
    "Une interface complète : 3 rôles, notifications, dashboard, RGPD",
    "Un système extensible : nouveau moteur IA = 1 fichier à swapper",
    "Une réduction estimée de 70% du temps de tri pour le recruteur",
], Inches(0.6), Inches(2.25), Inches(8.1), Inches(4.3),
   font_size=13, color=DARK_TEXT)

# Acquis techniques
add_rect(slide, Inches(9.2), Inches(1.7), Inches(3.7), Inches(5.1), WHITE)
add_rect(slide, Inches(9.2), Inches(1.7), Pt(5), Inches(5.1), BLUE)
add_text_box(slide, "🎓  Acquis du stage",
             Inches(9.4), Inches(1.8), Inches(3.4), Pt(28),
             font_size=14, bold=True, color=BLUE)
bullet_box(slide, [
    "Fine-tuning BERT / Transformers",
    "NLP Pipeline français",
    "FastAPI production-ready",
    "React SPA + Auth JWT",
    "PostgreSQL + Alembic ORM",
    "SMTP & OAuth 2.0",
    "Architecture microservices",
], Inches(9.4), Inches(2.25), Inches(3.3), Inches(4.3),
   font_size=13, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Perspectives
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, LIGHT_GRAY)
add_header_bar(slide, "Perspectives — 2 mois restants & evolutions futures")
add_footer(slide, 14)

# Court terme
add_rect(slide, Inches(0.4), Inches(1.7), Inches(12.5), Pt(26), RED_SOFT)
add_text_box(slide, "🔴  Court terme — à réaliser pour la soutenance finale",
             Inches(0.6), Inches(1.72), Inches(12), Pt(22),
             font_size=12, bold=True, color=WHITE)
bullet_box(slide, [
    "Notification email recruteur à chaque nouvelle candidature → boucle UX complète",
    "Planification d'entretiens : recruteur propose créneaux, candidat accepte → flux E2E complet",
    "Page candidat enrichie : score détaillé, skills manquants, conseils personnalisés",
], Inches(0.5), Inches(2.05), Inches(12.3), Inches(1.2),
   font_size=12.5, color=DARK_TEXT, bullet="  ›  ")

# Moyen terme
add_rect(slide, Inches(0.4), Inches(3.4), Inches(12.5), Pt(26), ACCENT)
add_text_box(slide, "🟡  Moyen terme — amélioration continue de l'IA",
             Inches(0.6), Inches(3.42), Inches(12), Pt(22),
             font_size=12, bold=True, color=WHITE)
bullet_box(slide, [
    "Feedback loop : recruteur note 'embauché / refusé' → données stockées → réentraînement BERT périodique",
    "Métriques d'évaluation : Precision@K, NDCG@K, MRR — mesure objective de la qualité du ranking",
    "Explication IA (XAI) : radar chart des dimensions, top skills manquantes par candidat",
], Inches(0.5), Inches(3.75), Inches(12.3), Inches(1.2),
   font_size=12.5, color=DARK_TEXT, bullet="  ›  ")

# Long terme
add_rect(slide, Inches(0.4), Inches(5.1), Inches(12.5), Pt(26), GREEN)
add_text_box(slide, "🟢  Long terme — production & évolutions",
             Inches(0.6), Inches(5.12), Inches(12), Pt(22),
             font_size=12, bold=True, color=WHITE)
bullet_box(slide, [
    "Docker Compose : backend + frontend + DB déployables en une commande → production-ready",
    "Tests automatisés CI/CD : pipeline de tests unitaires sur le matching (GitHub Actions)",
    "Application mobile (React Native) + API publique pour intégration RH tiers",
], Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.2),
   font_size=12.5, color=DARK_TEXT, bullet="  ›  ")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Merci / Questions
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, NAVY)
add_rect(slide, 0, Inches(5.5), W, Inches(2.0), BLUE)
add_rect(slide, 0, Inches(5.4), W, Pt(5), ACCENT)

add_logo(slide, LOGO_3S_PATH, Inches(0.5), Inches(0.35), Inches(1.4))
add_logo(slide, LOGO_ESPRIT_PATH, Inches(11.4), Inches(0.2), Inches(1.6))

add_text_box(slide, "Merci de votre attention",
             Inches(0), Inches(1.2), W, Inches(1.1),
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "Questions & Discussion",
             Inches(0), Inches(2.4), W, Inches(0.7),
             font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.5), Inches(3.2), Inches(4.3), Pt(4), ACCENT)

add_text_box(slide, "Youssef GARA  —  garayoussef0@gmail.com",
             Inches(0), Inches(3.5), W, Inches(0.5),
             font_size=16, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

add_text_box(slide, "3S TalentMatch  •  Matching IA CV/Offres  •  FastAPI + React + BERT v2.0",
             Inches(0), Inches(4.1), W, Inches(0.45),
             font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

# Technos bas de page
tech_badges = ["FastAPI", "React", "PostgreSQL", "BERT", "spaCy", "JWT", "SMTP"]
for i, tech in enumerate(tech_badges):
    x = Inches(1.2) + i * Inches(1.65)
    add_rect(slide, x, Inches(5.9), Inches(1.45), Inches(0.55), BLUE)
    add_text_box(slide, tech, x, Inches(5.95), Inches(1.45), Inches(0.42),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# NOTES ORATEUR — un paragraphe par slide
# ════════════════════════════════════════════════════════════════════════════

SPEAKER_NOTES = {

0: """Bonjour, je m'appelle Youssef Gara, etudiant en Licence Genie Logiciel a l'ESPRIT.
Je vais vous presenter aujourd'hui mon projet de stage de fin d'etudes realise au sein de la societe 3S, intitule TalentMatch.
Il s'agit d'une plateforme intelligente de recrutement basee sur l'intelligence artificielle.
Ce stage d'une duree de 6 mois, de fevrier a aout 2026, est encadre academiquement par Madame Loujayne Bouzrati et au sein de l'entreprise par Madame Sirine Nasri.
Je vous remercie de l'attention que vous portez a cette premiere restitution.""",

1: """La presentation se deroulera en sept parties.
Je commencerai par vous presenter le contexte et l'entreprise 3S, puis j'aborderai le cadre du projet en expliquant la problematique et l'etude de l'existant.
Ensuite viendra l'analyse des besoins, l'architecture technique retenue, la methodologie de developpement suivie sous forme de 6 sprints d'un mois.
Je vous montrerai ensuite l'etat d'avancement actuel du projet avec une demonstration du systeme.
Je terminerai par la conclusion et les perspectives pour les deux mois restants.""",

2: """3S est une entreprise tunisienne specialisee dans l'integration des infrastructures IT, fondee en 1988.
Avec plus de 35 ans d'expertise, elle compte aujourd'hui plus de 300 employes et plus de 450 certifications professionnelles cumulees.
Elle sert plus de 2000 clients dans des secteurs varies — aussi bien des grands groupes que des PME.
Elle a obtenu deux certifications internationales : l'ISO 9001 version 2015 pour la qualite de gestion, et l'ISO 27001 version 2013 pour la securite de l'information, ce qui en fait le premier integrateur IT tunisien certifie sur ce second referentiel.
En 2019, elle a etendu ses activites en ouvrant une filiale a Abidjan pour le marche ivoirien.
C'est dans ce contexte que ma mission de stage a ete definie : concevoir et developper une plateforme intelligente permettant d'automatiser le matching entre les CVs des candidats et les offres d'emploi.""",

3: """Le probleme que j'ai identifie est le suivant : chez 3S, comme dans beaucoup d'entreprises, le tri des CVs se fait encore manuellement.
Cela prend beaucoup de temps, introduit des biais humains, et ne passe pas a l'echelle quand le volume de candidatures augmente.
En etudiant les solutions existantes sur le marche, j'ai constate que des outils comme LinkedIn Recruiter ou Workday sont tres couteux et ne proposent pas un vrai traitement du langage naturel en francais.
Les systemes de tri par simples mots-cles manquent la semantique : un candidat qui ecrit "developpement logiciel" ne sera pas retrouve si l'offre dit "software development".
Notre solution, TalentMatch, repond a ces limites en proposant un matching semantique base sur BERT, adapte au marche tunisien, multilingue francais-anglais et entierement personnalisable.""",

4: """La solution que j'ai developpee est une plateforme web complete avec trois niveaux d'acces.
Le candidat peut s'inscrire, uploader son CV et suivre en temps reel l'etat de ses candidatures.
Le recruteur cree des offres d'emploi et recoit automatiquement les candidats classes par un score d'intelligence artificielle, avec le detail des competences matchees.
L'administrateur gere les utilisateurs, cree les comptes recruteurs et supervise les logs d'acces pour la conformite RGPD.
Tout le processus est automatise : des qu'un CV est uploade, il est analyse, parse par notre pipeline NLP, et un score de matching est calcule automatiquement par rapport a l'offre concernee.
Le recruteur n'a plus qu'a consulter le classement et prendre sa decision en un clic.""",

5: """En termes de besoins fonctionnels, j'ai identifie les fonctionnalites essentielles pour chaque acteur du systeme.
Pour le candidat : s'inscrire via email ou via OAuth Google et LinkedIn, uploader son CV dans les formats PDF, DOCX ou image, consulter les offres disponibles et postuler, puis recevoir des notifications sur l'evolution de sa candidature.
Pour le recruteur : creer et gerer les offres d'emploi, lancer le matching IA pour une offre donnee, consulter les candidats classes par score et modifier les statuts selon ses decisions.
Pour l'administrateur : gerer l'ensemble des comptes utilisateurs, creer des comptes recruteurs, consulter les logs d'acces pour la conformite RGPD et avoir une vue globale du systeme via un tableau de bord.""",

6: """Au-dela des fonctionnalites, le systeme doit respecter des contraintes non fonctionnelles importantes que j'ai organisees en six axes.
La securite : toutes les communications sont protegees par JWT avec refresh token, les mots de passe sont haches en bcrypt, et l'OAuth garantit une authentification fiable via des fournisseurs tiers.
La performance : le modele BERT est charge en memoire au demarrage du serveur pour eviter les delais — le score est calcule en moins de 2 secondes par candidat.
La conformite RGPD : chaque acces aux donnees est trace avec l'identite de l'utilisateur, l'action effectuee et l'horodatage, accessible uniquement par l'administrateur.
La maintenabilite : l'architecture est modulaire, les moteurs IA sont interchangeables et l'API est documentee automatiquement via Swagger.
Enfin, l'extensibilite permet d'ajouter un nouveau moteur IA en un seul fichier sans modifier le reste du systeme.""",

7: """L'architecture technique retenue est une architecture 3 tiers standard.
La couche frontend est developpee en React avec Vite : elle expose les interfaces des trois types d'utilisateurs et communique avec le backend via des appels HTTP REST securises par JWT.
La couche backend est developpee en FastAPI Python : elle expose l'ensemble des endpoints de l'API, gere la logique metier, l'authentification, et orchestre le pipeline d'intelligence artificielle.
La base de donnees est PostgreSQL, geree via SQLAlchemy comme ORM et Alembic pour le versionnage des migrations.
La couche IA repose sur notre modele TalentMatch-BERT v2.0, un modele multilingual MiniLM fine-tune sur des CVs reels, assiste de spaCy pour le parsing NLP en francais.
Ces quatre couches communiquent de maniere securisee via CORS configure et JWT Bearer token.""",

8: """Le coeur technique du projet est le pipeline de matching IA qui se deroulant en cinq etapes.
Premiere etape, l'extraction du texte brut : pour un PDF textuel on utilise PyMuPDF, pour un DOCX on utilise python-docx, et pour les PDFs scannes ou les images on fait appel a l'OCR Tesseract.
Deuxieme etape, le NLP Parser base sur spaCy analyse ce texte et extrait les competences, les experiences professionnelles, les formations et les informations de contact du candidat.
Troisieme etape, le modele TalentMatch-BERT v2.0 encode semantiquement le texte du CV et le texte de l'offre pour calculer une similarite vectorielle.
Quatrieme etape, un reseau de neurones MLP combine cinq dimensions : la similarite semantique, le taux de competences requises matchees, le score d'experience, le niveau de formation et le taux de competences appreciees.
Cinquieme etape, un score final de 0 a 100 pourcent est produit avec une decision automatique : Excellent candidat, Bon profil, A evaluer ou Non adapte.
Pour eviter les faux positifs, un systeme de plafonnement par domaine metier est applique : si le candidat est dans un domaine totalement different de l'offre, son score est plafonne automatiquement.""",

9: """Pour organiser le developpement sur 6 mois, j'ai adopte une methodologie Agile Scrum avec 6 sprints d'un mois chacun.
Le Sprint 1, en fevrier, a ete consacre a la phase de conception : cahier des charges, diagrammes UML, architecture technique et mise en place de l'environnement de developpement.
Le Sprint 2, en mars, a porte sur l'extraction des CVs : upload multi-format PDF et DOCX, integration de l'OCR EasyOCR pour les documents scannes, et stockage en base de donnees.
Le Sprint 3, en avril, a ete dedie au pipeline NLP avec spaCy pour extraire automatiquement les competences, les experiences et les formations a partir du texte brut des CVs.
Nous sommes actuellement dans le Sprint 4, en mai, qui integre le modele BERT, developpe le module de matching candidat-offre et construit le tableau de bord RH. C'est l'etat que je vous presente aujourd'hui dans cette premiere restitution.
Les Sprints 5 et 6, de juin a aout, couvriront la securite, les tests d'integration, les ameliorations UX, la documentation et la preparation de la soutenance finale.""",

10: """A ce stade de la premiere restitution, l'ensemble des fonctionnalites principales sont developpees et operationnelles.
Le systeme d'authentification est complet avec JWT, OAuth Google et LinkedIn, verification par email et trois niveaux de roles.
L'extraction de CVs fonctionne pour les trois formats : PDF textuel, DOCX, et PDF scanne via OCR.
Le NLP Parser analyse les CVs en francais et en anglais pour extraire competences, experiences et formations.
Le moteur de matching IA TalentMatch-BERT v2.0 est entraine et operationnel : il produit un score de 0 a 100 pourcent avec une decision automatique en quatre niveaux.
Le systeme de notifications envoie des emails aux candidats lors de l'acceptation ou du refus de leur candidature.
L'interface React est complete pour les trois profils utilisateurs : candidat, recruteur et administrateur.
Au total, le systeme a ete valide sur plus de 80 CVs reels en test, couvrant 14 domaines metiers differents.""",

11: """Je vais maintenant vous montrer le parcours type d'utilisation de la plateforme en quatre etapes.
Premiere etape : le candidat uploade son CV depuis son espace personnel. L'extraction du texte et le parsing NLP se font automatiquement en quelques secondes — les competences, experiences et formations sont detectees sans aucune intervention manuelle.
Deuxieme etape : le recruteur acceede a son offre et clique sur "Lancer le matching". Le systeme calcule les scores BERT pour tous les candidats ayant postule et les classe du meilleur au moins bon.
Troisieme etape : les resultats s'affichent avec pour chaque candidat son score en pourcentage, sa decision automatique — Excellent, Bon profil, A evaluer ou Non adapte — et le detail des competences matchees.
Quatrieme etape : le recruteur accepte ou refuse d'un clic. Le candidat recoit immediatement un email de notification et la cloche in-app se met a jour en temps reel.""",

12: """Pour conclure cette premiere restitution, TalentMatch repond concretement au besoin identifie chez 3S : automatiser et objectiver le processus de recrutement.
Le projet apporte un pipeline IA complet de bout en bout, depuis l'upload du CV jusqu'a la decision automatique, avec un modele BERT fine-tune sur des donnees reelles en francais et en anglais.
Sur le plan technique, ce stage m'a permis d'acquerir des competences solides en traitement du langage naturel, en fine-tuning de modeles Transformers, en developpement d'API REST avec FastAPI, et en developpement frontend React avec gestion d'etat et authentification.""",

13: """Pour les deux mois restants du stage, les priorites sont bien definies selon trois horizons.
A court terme, nous allons completer le flux de recrutement de bout en bout avec la planification des entretiens, ameliorer l'interface candidat pour lui afficher son score detaille et les competences manquantes, et finaliser les notifications email cote recruteur.
A moyen terme, nous mettrons en place une boucle de feedback : les decisions du recruteur — embauche ou refuse — seront stockees et utilisees pour re-entrainer le modele BERT periodiquement, afin d'ameliorer continuellement la qualite du matching sur les donnees reelles de 3S.
A long terme, le projet sera packague avec Docker Compose pour permettre un deploiement en production en une seule commande, et une suite de tests automatises sera mise en place pour garantir la maintenabilite du systeme sur la duree.""",

14: """Je vous remercie pour votre attention tout au long de cette presentation.
Je reste disponible pour repondre a vos questions concernant les choix techniques du projet, le fonctionnement du modele BERT, l'architecture du systeme ou les perspectives d'evolution pour la suite du stage.""",
}

# Injection des notes dans chaque slide
for slide_idx, note_text in SPEAKER_NOTES.items():
    try:
        s = prs.slides[slide_idx]
        tf = s.notes_slide.notes_text_frame
        tf.text = note_text
    except Exception:
        pass

# ════════════════════════════════════════════════════════════════════════════
# Sauvegarde
# ════════════════════════════════════════════════════════════════════════════
out_path = r"c:\Users\youssef\Desktop\3s-talentmatch\docs\Presentation_TalentMatch_Youssef_GARA_v7.pptx"
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"OK  Presentation generee : {out_path}")
print(f"    {len(prs.slides)} slides")
